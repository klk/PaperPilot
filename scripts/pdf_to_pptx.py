#!/usr/bin/env python3
"""Rebuild a PDF page as editable PowerPoint objects.

Text spans become editable text boxes, PDF rectangles and lines become PPT
shapes, and embedded PDF images remain independent image objects. This is a
best-effort layout converter; a commercial PDF engine is still preferable for
very complex transparency, clipping paths, or rotated text.
"""

from __future__ import annotations

import argparse
import io
import os
import platform
from typing import Any

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def clamp(value: float) -> int:
    return max(0, min(255, round(value * 255)))


def color(value: Any, fallback=(0, 0, 0)) -> RGBColor:
    if not value or not isinstance(value, (tuple, list)) or len(value) < 3:
        return RGBColor(*fallback)
    return RGBColor(clamp(value[0]), clamp(value[1]), clamp(value[2]))


def target_font() -> str:
    configured = os.environ.get("PDF2PPTX_FONT")
    if configured:
        return configured
    return "PingFang SC" if platform.system() == "Darwin" else "Noto Sans CJK SC"


def slide_xy(value: float, page_size: float, slide_size: int) -> int:
    return round(value / page_size * slide_size)


def add_pdf_drawing(slide, drawing, page_rect, slide_width, slide_height, scale):
    """Map the common PDF drawing primitives used by tables and borders."""
    page_width, page_height = page_rect.width, page_rect.height
    stroke = drawing.get("color")
    fill = drawing.get("fill")
    width = drawing.get("width") or 0.5
    for item in drawing.get("items", []):
        kind = item[0]
        if kind == "re":
            rect = item[1]
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                slide_xy(rect.x0, page_width, slide_width),
                slide_xy(rect.y0, page_height, slide_height),
                slide_xy(rect.width, page_width, slide_width),
                slide_xy(rect.height, page_height, slide_height),
            )
            if fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = color(fill)
            else:
                shape.fill.background()
            if stroke:
                shape.line.color.rgb = color(stroke)
                shape.line.width = Pt(max(0.25, float(width) * scale))
            else:
                shape.line.fill.background()
        elif kind == "l":
            start, end = item[1], item[2]
            line = slide.shapes.add_connector(
                1,
                slide_xy(start.x, page_width, slide_width),
                slide_xy(start.y, page_height, slide_height),
                slide_xy(end.x, page_width, slide_width),
                slide_xy(end.y, page_height, slide_height),
            )
            line.line.color.rgb = color(stroke)
            line.line.width = Pt(max(0.25, float(width) * scale))


def add_text_line(slide, line, page_rect, slide_width, slide_height, scale):
    spans = line.get("spans", [])
    if not spans:
        return
    bbox = fitz.Rect(line["bbox"])
    text = "".join(span.get("text", "") for span in spans)
    if not text.strip():
        return
    box = slide.shapes.add_textbox(
        slide_xy(bbox.x0, page_rect.width, slide_width),
        slide_xy(bbox.y0, page_rect.height, slide_height),
        max(slide_xy(max(1, bbox.width), page_rect.width, slide_width), 1),
        max(slide_xy(max(1, bbox.height), page_rect.height, slide_height), 1),
    )
    box.text_frame.clear()
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.word_wrap = False
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = box.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    for span in spans:
        run = paragraph.add_run()
        run.text = span.get("text", "")
        font = run.font
        font.name = target_font()
        font.size = Pt(max(1, float(span.get("size", 10)) * scale))
        font.color.rgb = color(span.get("color"), (0, 0, 0))
        font.bold = "bold" in str(span.get("font", "")).lower() or "semibold" in str(span.get("font", "")).lower()
    return box


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("input_pdf")
    parser.add_argument("output_pptx")
    args = parser.parse_args()

    pdf = fitz.open(args.input_pdf)
    if pdf.page_count == 0:
        raise ValueError("PDF 没有可转换的页面")

    first = pdf[0].rect
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = round(prs.slide_width * first.height / first.width)
    blank_layout = prs.slide_layouts[6]

    for page in pdf:
        page_rect = page.rect
        slide = prs.slides.add_slide(blank_layout)
        scale = prs.slide_width / (page_rect.width * 12700)

        for drawing in page.get_drawings():
            add_pdf_drawing(slide, drawing, page_rect, prs.slide_width, prs.slide_height, scale)

        for block in page.get_text("dict").get("blocks", []):
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                add_text_line(slide, line, page_rect, prs.slide_width, prs.slide_height, scale)

        # Keep uncommon PDF image blocks as independent editable picture
        # objects. They are not merged into a full-page background image.
        for block in page.get_text("dict").get("blocks", []):
            if block.get("type") != 1 or not block.get("image"):
                continue
            bbox = fitz.Rect(block["bbox"])
            slide.shapes.add_picture(
                io.BytesIO(block["image"]),
                slide_xy(bbox.x0, page_rect.width, prs.slide_width),
                slide_xy(bbox.y0, page_rect.height, prs.slide_height),
                width=slide_xy(bbox.width, page_rect.width, prs.slide_width),
                height=slide_xy(bbox.height, page_rect.height, prs.slide_height),
            )

    prs.save(args.output_pptx)
    pdf.close()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
